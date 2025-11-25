
import os.path

from typing import Any, Literal
from werkzeug.datastructures import FileStorage
from pptx.enum.shapes import MSO_SHAPE_TYPE as SHAPE_TYPE
from pptx.shapes.picture import Picture as PPTXPicture
from pptx import Presentation
from flask_login import current_user

from app.app import UPLOAD_PATH, hid
from app.models import Draft, DraftImage, DraftLabel, Set, Image, db

TEMP_UPLOAD_PATH = os.path.join(UPLOAD_PATH, "temp")

SlideItem = dict[str, str]
ItemsList = list[SlideItem]



def create_draft() -> int:
    """ Create a new draft and return its ID. """
    try:
        draft = Draft()
        db.session.add(draft)
        db.session.commit()
        return draft.id
    except Exception as e:
        db.session.rollback()
        raise e


def extract_images(presentation_file: FileStorage, draft_id: int) -> tuple[str, ItemsList, ItemsList] | Literal[False]:
    """
    Extract images and text labels from a presentation file and save them to the database.
    Returns a tuple containing the temporary presentation file address, a list of image filenames with slide numbers,
    and a list of text labels with slide numbers.
    """
    address = temp_save(presentation_file)
    if not address: return False

    target_directory = os.path.join(UPLOAD_PATH, "sets", f"draft_{draft_id}")

    pres = Presentation(address)
    draft = Draft.query.get(draft_id)
    if not isinstance(draft, Draft): return False
    pres_n = draft.presentations

    labels: list[dict[str, str]] = []
    images: list[dict[str, str]] = []
    _images: list[dict[str, Any]] = []

    image_n = 1
    for slide_n, slide in enumerate(pres.slides):
        for shape in slide.shapes:
            slide_encoded = hid.encode(pres_n * 10_000 + slide_n)
            if not isinstance(slide_encoded, str): slide_encoded = str(pres_n * 10_000 + slide_n)

            if shape.shape_type == SHAPE_TYPE.PICTURE:
                if not isinstance(shape, PPTXPicture): continue
                image = shape.image
                try:
                    filename = save_image(image.blob, image.ext, target_directory, start=image_n)
                except Exception:
                    filename = False
                if not filename: continue
                print(f"Saved image {filename} from slide {slide_n + 1}")
                draft_img = DraftImage(draft_id, filename, pres_n, slide_n, label = "")
                db.session.add(draft_img)
                _images.append({"id": draft_img, "filename": filename, "label": "", "slide": slide_encoded})
                image_n += 1
            elif shape.has_text_frame:
                text = getattr(shape, "text", "")
                if not text: continue
                l = save_labels(text, pres_n,   slide_n, draft_id)
                for s in l: labels.append({"text": s, "slide": slide_encoded})
    db.session.commit()
    for img in _images:
        images.append({"id": img["id"].id, "filename": img["filename"], "label": img["label"], "slide": img["slide"]})
    
    return address, images, labels


def temp_save(file: FileStorage) -> str | Literal[False]:
    """ Save a file temporarily and return its path. """
    if file.filename is None: return False

    EXTENSION = file.filename.rsplit('.', 1)[-1].lower()
    FILENAME = get_free_filename(TEMP_UPLOAD_PATH, EXTENSION)
    path = os.path.join(TEMP_UPLOAD_PATH, FILENAME)
    os.makedirs(TEMP_UPLOAD_PATH, exist_ok = True)

    file.save(path)
    return path


def temp_remove(file_path: str) -> bool:
    """ Remove a temporary file. """
    if not os.path.exists(file_path): return False
    os.remove(file_path)
    return True


def save_image(image_bytes: bytes, ext: str, target_directory: str, start: int = 1) -> str | Literal[False]:
    """ Save image bytes to a file in the target directory and return the filename. """
    if not image_bytes: return False
    if ext.lower() not in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff']: return False

    os.makedirs(target_directory, exist_ok = True)

    filename = get_free_filename(target_directory, ext, "img", start)
    path = os.path.join(target_directory, filename)

    with open(path, 'wb') as f:
        f.write(image_bytes)

    return filename


def get_free_filename(dir: str, ext: str, prefix: str = "tmp", start: int = 1) -> str:
    """ Generate a free filename in the specified directory with the given extension and prefix. """
    counter = start - 1
    while True:
        counter += 1
        filename = f"{prefix}_{counter:0>6}.{ext}"
        path = os.path.join(dir, filename)
        if not os.path.exists(path): return filename


def get_free_index(dir: str, prefix: str = "tmp", ext: str = "*", start: int = 1) -> int:
    """ Get the next free index for filenames in the specified directory with the given prefix and extension. """
    counter = start
    os.makedirs(dir, exist_ok = True)
    while True:
        # if not os.path.exists(path): return counter
        for f in os.listdir(dir):
            if f.startswith(f"{prefix}_{counter:0>6}.") and (ext == "*" or f.endswith(f".{ext}")): break
        else: return counter
        counter += 1


def save_labels(text: str, pres_n: int, slide: int, draft_id: int) -> list[str]:
    """ Save text labels to the database and return a list of labels. """
    if not text: return []
    labels = [line.strip() for line in text.splitlines() if line.strip()]
    for label in labels:
        draft_label = DraftLabel(draft_id, label, pres_n, slide)
        db.session.add(draft_label)
    return labels