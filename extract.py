
from pptx import Presentation
from pptx.slide import Slide
from pptx.parts.image import Image
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE


from typing import Generator
from werkzeug.utils import secure_filename
from easygui import fileopenbox
import os



DATA_FOLDER = os.path.join(os.getcwd(), "data")


# filename = fileopenbox(default = os.path.join(DATA_FOLDER, "*.pptx"),
#                         filetypes = ["*.zip", "*.pptx"],
#                         title = "Vyberte soubor")

filename = os.path.join(DATA_FOLDER, "Egypt.pptx")

if filename is None:
    raise FileNotFoundError("You have to choose a file")

if not isinstance(filename, str):
    raise FileNotFoundError("Something went wrong: Failed to open the file")

print(f"{filename=}\n{DATA_FOLDER=}\n")



def iter_picture_shapes(slide: Slide) -> Generator[BaseShape, None, None]:
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape



prs = Presentation(filename)

last_slide:list[list[bytes]] = []

for slide_n, slide in enumerate(prs.slides):
    last_slide.append([])

    for img_n, picture in enumerate(iter_picture_shapes(slide)):
        image = picture.image

        if not isinstance(image, Image): raise TypeError("Not an image")
        image_bytes = image.blob

        if len(last_slide) > 0 and image_bytes in last_slide[-1]: continue  # Already processed on this slide
        
        last_slide[-1].append(image_bytes)
        if len(last_slide) > 1 and image_bytes in last_slide[-2]: continue  # Already processed on the previous slide


        filename = f"{slide_n}_{img_n}.{image.ext}"

        image_path = os.path.join(DATA_FOLDER, secure_filename(filename))

        with open(image_path, 'wb') as f:
            f.write(image_bytes)

    if slide_n >= 2:
        del last_slide[-3]

