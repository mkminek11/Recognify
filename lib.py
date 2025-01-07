from werkzeug.datastructures import FileStorage
from werkzeug.utils import secure_filename
from typing import Generator, Literal
import os
import base64

from pptx import Presentation as PPTX
from pptx.slide import Slide as PPTXSlide
from pptx.parts.image import Image as PPTXImage
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture as PPTXPicture
from pptx.enum.shapes import MSO_SHAPE_TYPE

from db import Presentation, Image, Label, app, db


def _iter_shapes(slide: PPTXSlide) -> Generator[BaseShape, None, None]:
    """
    Yields all picture shapes in a slide.

    Args:
        slide: `Slide` - The slide to yield shapes from.

    Yields:
        `BaseShape` - Shapes of type `MSO_SHAPE_TYPE.PICTURE`.
    """

    for shape in slide.shapes:
        yield shape



def upload_file(root_path: str, file: FileStorage, uuid: str) -> str:
    """
    Creates a subdirectory in `user_upload` and saves it under a unique name (uuid4).
    The file is then saved in the subdirectory.

    Args:
        root_path: `str` - The root path of the application.
        file: `FileStorage` - The file to be saved.

    Returns:
        `str` - The unique id of the presentation, also the name of the subdirectory.
    """

    uploads_dir = os.path.join(root_path, "user_upload")
    base_dir_path = os.path.join(uploads_dir, uuid)
    os.makedirs(base_dir_path, exist_ok = True)
    # file.save(os.path.join(base_dir_path, "presentation.pptx"))


def extract_images(pres_file: FileStorage, pres_uuid: str) -> bool:
    """
    Extracts images from a presentation and saves them to a directory.
    The images are saved in a subdirectory named after the uuid of the presentation.
    The images are named after the image order in the presentation.

    Args:
        presentation: `FileStorage` - The presentation file to extract images from.
        pres_uuid: `str` - The uuid of the presentation.

    Returns:
        `bool` - Whether the extraction was successful.
    """

    pres_file_name = pres_file.filename
    if not pres_file_name or not pres_file_name.endswith(".pptx"): return False
    pres_title = os.path.splitext(os.path.basename(pres_file_name))[0].capitalize()

    pres_obj = Presentation(uuid = pres_uuid, title = pres_title)
    db.session.add(pres_obj)
    db.session.commit()

    pres_id = pres_obj.id

    root_dir = os.path.join(app.root_path, "user_upload", pres_uuid)
    images_dir = os.path.join(root_dir, "images")
    os.makedirs(images_dir, exist_ok = True)

    pptx_slides = PPTX(os.path.join(app.root_path, "user_upload", pres_uuid, "presentation.pptx")).slides

    last_slide:list[list[bytes]] = []
    image_index = 0

    for slide_n, slide in enumerate(pptx_slides):
        last_slide.append([])

        for img_n, shape in enumerate(_iter_shapes(slide)):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if not isinstance(shape, PPTXPicture): continue
                image_path = _save_image(shape.image, last_slide, image_index, images_dir)
                image_index += 1
                if not image_path: continue

                db.session.add(Image(file = image_path, slide = slide_n, presentation = pres_id))

            elif shape.has_text_frame:
                if not hasattr(shape, "text"): continue
                if not shape.text: continue

                db.session.add(Label(text = shape.text, slide = slide_n, presentation = pres_id))

        if slide_n >= 2:
            del last_slide[-3]

    db.session.commit()
    return True



def _save_image(image: PPTXImage, last_slide: list[list[bytes]], img_n: int, target_dir: str) -> str | Literal[False]:
    """
    Saves an image to the target directory, but only if it has not been already saved in the current slide or the previous one.

    Args:
        image: `PPTXImage` - The image to be saved.
        last_slide: `list[list[bytes]]` - A list of slides, where each sublist contains the image bytes of all images
            in one slide. The most recent slide is at the end of the list.
        img_n: `int` - The index of the image in the current slide.
        target_dir: `str` - The directory where the image should be saved.

    Returns:
        `str | Literal[False]` - The path of the saved image, or `False` if the image was already saved in the current or
            previous slide.
    """

    if not isinstance(image, PPTXImage): return False
    image_bytes = image.blob

    if len(last_slide) > 0 and image_bytes in last_slide[-1]: return False  # Already processed on this slide
    
    last_slide[-1].append(image_bytes)
    if len(last_slide) > 1 and image_bytes in last_slide[-2]: return False  # Already processed on the previous slide

    try:
        file_name = f"{img_n:03}.{image.ext}"
    except ValueError:
        print("Failed to save image")
        return False

    image_path = os.path.join(target_dir, secure_filename(file_name))

    with open(image_path, 'wb') as f:
        f.write(image_bytes)

    return image_path



def get_image_data(image_path: str, presentation_uuid: str) -> str:
    """
    Retrieves and encodes image data as a base64 string.

    Args:
        image_name: `str` - The name of the image file.
        presentation_uuid: `str` - The unique identifier of the presentation.

    Returns:
        `str` - The base64 encoded string of the image data.
    """

    with open(image_path, "rb") as f:
        return base64.b64encode(f.read()).decode('ascii')